"""
pptx_builder.py
Builds the final report PPTX by filling the template.pptx with campaign data.

Strategy:
  - Open template.pptx
  - Slides 1-10 are TEMPLATE slides (cover, top3 ttok, top3 xhs, sample posting, etc.)
  - We DUPLICATE template slides for each KOC/KOL and fill them
  - We delete unused template "sample" slides at the end

Template slide indices (0-based):
  0: Cover
  1: Top 3 TikTok
  2: Top 3 XHS
  3: TIKTOK KOL Posting (sample — duplicated per KOL)
  4: XHS KOC Posting (sample — duplicated per KOC)
  5: IG (EXTRA) divider
  6: TIKTOK IG Posting (sample — duplicated)
  7: XHS IG Posting (sample — duplicated)
  8: Campaign Summary
  9: Thank You
"""

import copy
import os
import re
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
from lxml import etree

# =============================================================
# Constants
# =============================================================
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# CHANGED: default to Watsons teal instead of generic green
DEFAULT_HEADER_COLOR = (0, 160, 168)
WHITE = RGBColor(255, 255, 255)
MAX_IMAGE_LONG_EDGE = 1400
JPEG_QUALITY = 78


# =============================================================
# URL cleaning
# =============================================================
def clean_url(text: str) -> str:
    """
    Extract just the URL from a string that might have extra share text.
    XHS often gives '逛到脚软❗️Watsons这波也太好🛍️了 http://xhslink.com/o/Uh9qVjaBXa  Copy and open rednote to view the note'
    We want to keep only 'http://xhslink.com/o/Uh9qVjaBXa'.
    """
    if not text:
        return ""
    # Find the first http(s) URL in the string
    m = re.search(r'(https?://[^\s\u4e00-\u9fff]+)', text)
    if m:
        return m.group(1).rstrip(',.;:!?)')
    return text.strip()


def prepare_image_for_ppt(image_path: str) -> str:
    """
    Compress screenshot-like images before inserting them into PPT.
    Keeps transparent PNGs unchanged so client logos stay clean.
    """
    if not image_path or not os.path.exists(image_path):
        return image_path

    try:
        with Image.open(image_path) as img:
            has_alpha = img.mode in ("RGBA", "LA") or (
                img.mode == "P" and "transparency" in img.info
            )
            if has_alpha:
                return image_path

            img = img.convert("RGB")
            w, h = img.size
            long_edge = max(w, h)

            if long_edge > MAX_IMAGE_LONG_EDGE:
                scale = MAX_IMAGE_LONG_EDGE / long_edge
                new_size = (max(1, int(w * scale)), max(1, int(h * scale)))
                img = img.resize(new_size, Image.Resampling.LANCZOS)

            fd, out_path = tempfile.mkstemp(prefix="ppt_img_", suffix=".jpg")
            os.close(fd)
            img.save(out_path, format="JPEG", quality=JPEG_QUALITY, optimize=True)

        if os.path.getsize(out_path) < os.path.getsize(image_path):
            return out_path
        try:
            os.remove(out_path)
        except OSError:
            pass
        return image_path
    except Exception:
        return image_path

# Template slide index constants
IDX_COVER = 0
IDX_TOP3_TIKTOK = 1
IDX_TOP3_XHS = 2
IDX_TIKTOK_POSTING = 3
IDX_XHS_POSTING = 4
IDX_IG_DIVIDER = 5
IDX_TIKTOK_IG = 6
IDX_XHS_IG = 7
IDX_SUMMARY = 8
IDX_THANKYOU = 9


# =============================================================
# Helpers — slide manipulation
# =============================================================
def duplicate_slide(prs, slide):
    """
    Duplicates a slide and appends it to the end of the presentation.
    Returns the new slide.
    """
    blank = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank)

    # Copy non-layout relationships first. Duplicated picture XML keeps the
    # original rIds, so we need to create matching relationships on the new
    # slide and then rewrite the copied XML to those new rIds.
    rel_id_map = {}
    for rId, rel in slide.part.rels.items():
        if rel.reltype.endswith("/slideLayout"):
            continue
        if rel.is_external:
            new_rId = new_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
        rel_id_map[rId] = new_rId

    # Copy all shapes from source slide
    for shape in slide.shapes:
        el = shape.element
        new_el = copy.deepcopy(el)
        for element in new_el.iter():
            for attr_name, attr_value in list(element.attrib.items()):
                if attr_value in rel_id_map:
                    element.attrib[attr_name] = rel_id_map[attr_value]
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    force_white_background(new_slide)

    return new_slide


def force_white_background(slide):
    """Force an explicit white slide background for PowerPoint dark mode."""
    try:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE
    except Exception:
        pass


def force_all_slides_white(prs):
    """Force every slide in the deck to use an explicit white background."""
    for slide in prs.slides:
        force_white_background(slide)


def reorder_slides(prs, new_order_indices):
    """
    Reorders slides in the presentation.
    new_order_indices: list of current slide indices in desired final order.
    """
    sldIdLst = prs.slides._sldIdLst
    slides = list(sldIdLst)
    if len(new_order_indices) != len(slides):
        raise ValueError(f"Order length {len(new_order_indices)} != slide count {len(slides)}")
    for el in slides:
        sldIdLst.remove(el)
    for i in new_order_indices:
        sldIdLst.append(slides[i])


def delete_slide(prs, slide_idx):
    """Delete slide at index."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[slide_idx])


# =============================================================
# Helpers — replace placeholders in a slide
# =============================================================
def replace_text_in_slide(slide, replacements: dict):
    """
    Replace {PLACEHOLDER} text in all shapes of a slide.
    Preserves formatting (font, size, color).
    """
    for shape in slide.shapes:
        # Regular text frames
        if shape.has_text_frame:
            _replace_in_text_frame(shape.text_frame, replacements)
        # Tables
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    _replace_in_text_frame(cell.text_frame, replacements)


def _replace_in_text_frame(tf, replacements: dict):
    for para in tf.paragraphs:
        for run in para.runs:
            for k, v in replacements.items():
                if k in run.text:
                    run.text = run.text.replace(k, str(v))


# =============================================================
# Helpers — replace a placeholder shape with an image
# =============================================================
def replace_shape_with_image(slide, shape_name: str, image_path: str, *,
                             keep_aspect_ratio: bool = True,
                             horizontal_align: str = "center"):
    """
    Find a shape by name, get its position+size, delete it, add an image in its place.
    If keep_aspect_ratio=True, fits image inside the box (with letterboxing if needed).
    """
    target = None
    for shape in slide.shapes:
        if shape.name == shape_name:
            target = shape
            break

    if target is None:
        return False

    x, y, w, h = target.left, target.top, target.width, target.height

    # Remove the placeholder shape
    sp = target._element
    sp.getparent().remove(sp)

    # Also remove any text label inside this shape area (e.g. "{HERO_IMAGE}")
    # We do this by checking textboxes that overlap significantly with this area
    # — simpler approach: just continue, text labels get covered by image

    if not image_path:
        return False

    try:
        image_path = prepare_image_for_ppt(image_path)
        if keep_aspect_ratio:
            # Calculate aspect-fit dimensions
            with Image.open(image_path) as img:
                iw, ih = img.size
            box_w_emu = w
            box_h_emu = h
            box_aspect = box_w_emu / box_h_emu
            img_aspect = iw / ih

            if img_aspect > box_aspect:
                # Image is wider — fit by width, center vertically
                new_w = box_w_emu
                new_h = int(box_w_emu / img_aspect)
                new_x = x
                new_y = y + (box_h_emu - new_h) // 2
            else:
                # Image is taller — fit by height
                new_h = box_h_emu
                new_w = int(box_h_emu * img_aspect)
                if horizontal_align == "left":
                    new_x = x
                elif horizontal_align == "right":
                    new_x = x + (box_w_emu - new_w)
                else:
                    new_x = x + (box_w_emu - new_w) // 2
                new_y = y

            slide.shapes.add_picture(image_path, new_x, new_y, width=new_w, height=new_h)
        else:
            slide.shapes.add_picture(image_path, x, y, width=w, height=h)
        return True
    except Exception as e:
        print(f"Image replacement error ({shape_name}): {e}")
        return False


def remove_text_shape_by_text(slide, text_match: str):
    """Remove any text-containing shape whose text contains the given match string."""
    to_remove = []
    for shape in slide.shapes:
        if shape.has_text_frame and text_match in shape.text_frame.text:
            to_remove.append(shape)
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)


def set_text_shape_hyperlink(slide, shape_name: str, url: str):
    """Make the text in a named shape clickable."""
    clean = clean_url(url)
    if not clean.startswith(("http://", "https://")):
        return False

    for shape in slide.shapes:
        if shape.name != shape_name or not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if not para.runs:
                run = para.add_run()
                run.text = clean
            for run in para.runs:
                if run.text.strip():
                    run.hyperlink.address = clean
        return True
    return False


def align_text_shape_to_shape(slide, text_shape_name: str, anchor_shape_name: str):
    """Align a text box with another shape's left and width."""
    text_shape = None
    anchor_shape = None
    for shape in slide.shapes:
        if shape.name == text_shape_name:
            text_shape = shape
        elif shape.name == anchor_shape_name:
            anchor_shape = shape

    if text_shape is None or anchor_shape is None:
        return False

    text_shape.left = anchor_shape.left
    text_shape.width = anchor_shape.width
    if text_shape.has_text_frame:
        for para in text_shape.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
    return True


# =============================================================
# Helpers — change header bar color
# =============================================================
def update_header_bar_color(slide, rgb: tuple):
    """Change the fill color of the 'ph_header_bar' or 'ph_divider_bar' shape."""
    r, g, b = rgb
    color = RGBColor(r, g, b)
    for shape in slide.shapes:
        if shape.name in ("ph_header_bar", "ph_divider_bar"):
            try:
                shape.fill.solid()
                shape.fill.fore_color.rgb = color
            except Exception:
                pass


def update_campaign_name_color(slide, rgb: tuple):
    """Change the campaign name color on cover slide."""
    r, g, b = rgb
    color = RGBColor(r, g, b)
    for shape in slide.shapes:
        if shape.name == "ph_campaign_name":
            try:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = color
            except Exception:
                pass


def update_thankyou_accent_color(slide, rgb: tuple):
    """Update accent lines and small headings on Thank You slide to match header color."""
    r, g, b = rgb
    color = RGBColor(r, g, b)
    for shape in slide.shapes:
        # Accent lines (top/bottom)
        if shape.name in ("ph_accent_top", "ph_accent_bottom"):
            try:
                shape.fill.solid()
                shape.fill.fore_color.rgb = color
            except Exception:
                pass
        # Accent labels (FOLLOW / VISIT / ENQUIRIES)
        if shape.name in ("ph_accent_label_follow", "ph_accent_label_visit", "ph_accent_label_enquiries"):
            try:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = color
            except Exception:
                pass


# =============================================================
# Number formatting
# =============================================================
def fmt(n):
    """Format a number with thousands separators."""
    try:
        return f"{int(n):,}"
    except (ValueError, TypeError):
        return str(n)


def engagement(item):
    """Compute engagement = likes + comments + saves + shares."""
    return item.get("likes", 0) + item.get("comments", 0) + item.get("saves", 0) + item.get("shares", 0)


def total(arr, key):
    return sum(item.get(key, 0) for item in arr)


# =============================================================
# Main builder
# =============================================================
def build_report(
    *,
    template_path: str,
    output_path: str,
    campaign_name: str,
    campaign_month: str,
    client_logo_path: str,
    hero_image_path: str = None,
    header_color: tuple = (124, 179, 66),
    campaign_data: dict,
):
    """
    Build the final PPT report.
    campaign_data structure (after drive_reader + image_downloader + ocr_classifier):
      {
        'xhs_kocs': [{'creator', 'post_url', 'views', 'likes', ...,
                      'post_screenshot_local', 'best_insight_local',
                      'ig_creator', 'ig_post_screenshot_local', 'best_ig_insight_local', ig_views, ...}, ...],
        'tiktok_kols': [...]
      }
    """
    prs = Presentation(template_path)

    xhs = campaign_data.get("xhs_kocs", [])
    tiktok = campaign_data.get("tiktok_kols", [])

    # ============================================================
    # SLIDE 1 — COVER
    # ============================================================
    cover = prs.slides[IDX_COVER]
    replace_text_in_slide(cover, {
        "{CAMPAIGN_NAME}": campaign_name.upper(),
        "{MONTH}": campaign_month,
    })
    update_campaign_name_color(cover, header_color)

    # Replace client logo placeholder
    if client_logo_path:
        replace_shape_with_image(cover, "ph_client_logo", client_logo_path)
        remove_text_shape_by_text(cover, "{CLIENT_LOGO}")

    # Replace hero image placeholder (if provided)
    if hero_image_path:
        replace_shape_with_image(cover, "ph_hero_image", hero_image_path)
        remove_text_shape_by_text(cover, "{HERO_IMAGE}")
    else:
        # No hero — remove placeholder shape and its text label
        for shape in list(cover.shapes):
            if shape.name == "ph_hero_image":
                sp = shape._element
                sp.getparent().remove(sp)
                break
        remove_text_shape_by_text(cover, "{HERO_IMAGE}")

    # ============================================================
    # SLIDE 2 — TOP 3 TIKTOK
    # ============================================================
    top3_tt = sorted(tiktok, key=lambda x: -x.get("views", 0))[:3]
    s = prs.slides[IDX_TOP3_TIKTOK]
    update_header_bar_color(s, header_color)
    if client_logo_path:
        replace_shape_with_image(s, "ph_client_logo", client_logo_path)
        remove_text_shape_by_text(s, "{CLIENT_LOGO}")
    _fill_top3(s, top3_tt, platform_label="TIKTOK")

    # ============================================================
    # SLIDE 3 — TOP 3 XHS
    # ============================================================
    top3_xhs = sorted(xhs, key=lambda x: -x.get("views", 0))[:3]
    s = prs.slides[IDX_TOP3_XHS]
    update_header_bar_color(s, header_color)
    if client_logo_path:
        replace_shape_with_image(s, "ph_client_logo", client_logo_path)
        remove_text_shape_by_text(s, "{CLIENT_LOGO}")
    _fill_top3(s, top3_xhs, platform_label="XHS")

    # ============================================================
    # POSTING SLIDES — duplicate the template per KOC/KOL
    # ============================================================
    # Build all the new posting slides at the end of the deck, then reorder.

    tt_posting_template = prs.slides[IDX_TIKTOK_POSTING]
    xhs_posting_template = prs.slides[IDX_XHS_POSTING]

    new_tt_posting_slides = []
    for kol in tiktok:
        new_slide = duplicate_slide(prs, tt_posting_template)
        _fill_posting_slide(
            new_slide, kol, header_title="TIKTOK KOL POSTING",
            client_logo_path=client_logo_path, header_color=header_color,
            use_ig_data=False
        )
        new_tt_posting_slides.append(new_slide)

    new_xhs_posting_slides = []
    for koc in xhs:
        new_slide = duplicate_slide(prs, xhs_posting_template)
        _fill_posting_slide(
            new_slide, koc, header_title="XHS KOC POSTING",
            client_logo_path=client_logo_path, header_color=header_color,
            use_ig_data=False
        )
        new_xhs_posting_slides.append(new_slide)

    # IG (Extra) postings
    tt_ig_template = prs.slides[IDX_TIKTOK_IG]
    xhs_ig_template = prs.slides[IDX_XHS_IG]

    new_tt_ig_slides = []
    for kol in tiktok:
        if not kol.get("ig_post_url"):
            continue
        new_slide = duplicate_slide(prs, tt_ig_template)
        _fill_posting_slide(
            new_slide, kol, header_title="TIKTOK KOL EXTRA POSTING (IG)",
            client_logo_path=client_logo_path, header_color=header_color,
            use_ig_data=True
        )
        new_tt_ig_slides.append(new_slide)

    new_xhs_ig_slides = []
    for koc in xhs:
        if not koc.get("ig_post_url"):
            continue
        new_slide = duplicate_slide(prs, xhs_ig_template)
        _fill_posting_slide(
            new_slide, koc, header_title="XHS KOC EXTRA POSTING (IG)",
            client_logo_path=client_logo_path, header_color=header_color,
            use_ig_data=True
        )
        new_xhs_ig_slides.append(new_slide)

    # ============================================================
    # IG DIVIDER
    # ============================================================
    div_slide = prs.slides[IDX_IG_DIVIDER]
    update_header_bar_color(div_slide, header_color)
    if client_logo_path:
        replace_shape_with_image(div_slide, "ph_client_logo", client_logo_path)
        remove_text_shape_by_text(div_slide, "{CLIENT_LOGO}")

    # ============================================================
    # SUMMARY
    # ============================================================
    summary_slide = prs.slides[IDX_SUMMARY]
    update_header_bar_color(summary_slide, header_color)
    if client_logo_path:
        replace_shape_with_image(summary_slide, "ph_client_logo", client_logo_path)
        remove_text_shape_by_text(summary_slide, "{CLIENT_LOGO}")
    _fill_summary(summary_slide, campaign_data, campaign_name, campaign_month)

    # ============================================================
    # THANK YOU (update accent color)
    # ============================================================
    thankyou_slide = prs.slides[IDX_THANKYOU]
    update_thankyou_accent_color(thankyou_slide, header_color)

    # ============================================================
    # REORDER: put new posting slides in the right place, delete templates
    # ============================================================
    # Current order (after additions):
    #   0..9: original templates
    #   10..: new slides we duplicated
    # We want final order:
    #   cover
    #   top3 tiktok
    #   top3 xhs
    #   [tiktok posting 1..N]
    #   [xhs posting 1..N]
    #   ig divider
    #   [tiktok ig 1..N]
    #   [xhs ig 1..N]
    #   summary
    #   thank you

    total_slides = len(prs.slides)
    n_template = 10  # original template count

    # Calculate index ranges for the NEW (duplicated) slides
    # They were appended in this order: tt_posting, xhs_posting, tt_ig, xhs_ig
    cursor = n_template
    tt_post_range = list(range(cursor, cursor + len(new_tt_posting_slides))); cursor += len(new_tt_posting_slides)
    xhs_post_range = list(range(cursor, cursor + len(new_xhs_posting_slides))); cursor += len(new_xhs_posting_slides)
    tt_ig_range = list(range(cursor, cursor + len(new_tt_ig_slides))); cursor += len(new_tt_ig_slides)
    xhs_ig_range = list(range(cursor, cursor + len(new_xhs_ig_slides))); cursor += len(new_xhs_ig_slides)

    new_order = (
        [IDX_COVER, IDX_TOP3_TIKTOK, IDX_TOP3_XHS]
        + tt_post_range
        + xhs_post_range
        + [IDX_IG_DIVIDER]
        + tt_ig_range
        + xhs_ig_range
        + [IDX_SUMMARY, IDX_THANKYOU]
    )

    # We need to also "skip" the template posting slides (3, 4, 6, 7) — they're in the original 10
    # but we've replaced them with duplicates. So they shouldn't appear in new_order.
    # Verify: new_order doesn't include 3, 4, 6, 7
    assert IDX_TIKTOK_POSTING not in new_order
    assert IDX_XHS_POSTING not in new_order
    assert IDX_TIKTOK_IG not in new_order
    assert IDX_XHS_IG not in new_order

    # Now we need a tricky bit: reorder all current slides, then delete the unused templates.
    # Strategy: build the final list of indices to KEEP in order, delete the rest by index from the back.

    # First, reorder all slides — but reorder requires same length. So we include templates we'll delete at the end.
    # Order: [our desired slides] + [templates we want to delete]
    leftover = [IDX_TIKTOK_POSTING, IDX_XHS_POSTING, IDX_TIKTOK_IG, IDX_XHS_IG]
    full_order = new_order + leftover
    assert len(full_order) == total_slides, f"Order length {len(full_order)} != total {total_slides}"

    reorder_slides(prs, full_order)

    # Now the last 4 slides are the unused templates — delete them
    for _ in range(len(leftover)):
        delete_slide(prs, len(prs.slides) - 1)

    # ============================================================
    # SAVE
    # ============================================================
    force_all_slides_white(prs)
    prs.save(output_path)
    return output_path


# =============================================================
# Sub-helpers for filling specific slides
# =============================================================
def _fill_top3(slide, top3_list, platform_label: str):
    """Fill a Top 3 slide with the top 3 KOC/KOLs."""
    while len(top3_list) < 3:
        top3_list = list(top3_list) + [None]

    for i in range(3):
        item = top3_list[i]
        idx = i + 1

        if item is None:
            # Hide this card
            replace_text_in_slide(slide, {
                f"{{TOP3_NAME_{idx}}}": "—",
                f"{{TOP3_FULLNAME_{idx}}}": "",
                f"{{TOP3_VIEWS_{idx}}}": "",
                f"{{TOP3_ENG_{idx}}}": "",
            })
            continue

        eng = engagement(item)

        # Replace text first
        # Truncate name on placeholder card if too long
        display_name = item["creator"]
        if len(display_name) > 14:
            display_name = display_name[:12] + "…"

        replace_text_in_slide(slide, {
            f"{{TOP3_NAME_{idx}}}": f"{display_name} · {platform_label}",
            f"{{TOP3_FULLNAME_{idx}}}": item["creator"],
            f"{{TOP3_VIEWS_{idx}}}": fmt(item["views"]),
            f"{{TOP3_ENG_{idx}}}": fmt(eng),
        })

        # Replace the colored card with the post screenshot
        if item.get("post_screenshot_local"):
            replace_shape_with_image(slide, f"ph_top3_image_{idx}",
                                     item["post_screenshot_local"])
            # Remove the "{name} · platform" label that was inside the colored card
            # (since it's now covered by the image, but we should also remove the text shape)
            for shape in list(slide.shapes):
                if shape.name == f"ph_top3_label_{idx}":
                    sp = shape._element
                    sp.getparent().remove(sp)
                    break


def _fill_posting_slide(slide, item, *, header_title: str, client_logo_path: str,
                       header_color: tuple, use_ig_data: bool = False):
    """Fill an individual posting slide."""
    # Update header
    update_header_bar_color(slide, header_color)
    replace_text_in_slide(slide, {"{POSTING_TYPE}": header_title})

    # Replace header text — find the textbox over the bar that has the original title
    for shape in slide.shapes:
        if shape.name == "ph_header_title":
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = header_title
            break

    # Replace client logo
    if client_logo_path:
        replace_shape_with_image(slide, "ph_client_logo", client_logo_path)
        remove_text_shape_by_text(slide, "{CLIENT_LOGO}")

    # Get fields based on whether we use IG data or main platform data
    if use_ig_data:
        creator = item.get("ig_creator") or item.get("creator")
        post_url = item.get("ig_post_url", "")
        post_ss = item.get("ig_post_screenshot_local")
        insight_ss = item.get("best_ig_insight_local")
        views = item.get("ig_views", 0)
        likes = item.get("ig_likes", 0)
        comments = item.get("ig_comments", 0)
        saves = item.get("ig_saves", 0)
        shares = item.get("ig_shares", 0)
    else:
        creator = item.get("creator", "")
        post_url = item.get("post_url", "")
        post_ss = item.get("post_screenshot_local")
        insight_ss = item.get("best_insight_local")
        views = item.get("views", 0)
        likes = item.get("likes", 0)
        comments = item.get("comments", 0)
        saves = item.get("saves", 0)
        shares = item.get("shares", 0)

    # Replace text placeholders
    replace_text_in_slide(slide, {
        "{CREATOR_NAME}": creator,
        "{POST_URL}": "Posted Link" if clean_url(post_url) else "",
        "{LIKES}": fmt(likes),
        "{COMMENTS}": fmt(comments),
        "{VIEWS}": fmt(views),
        "{SAVES}": fmt(saves),
        "{SHARES}": fmt(shares),
    })
    align_text_shape_to_shape(slide, "ph_post_url", "ph_insight_image")
    set_text_shape_hyperlink(slide, "ph_post_url", post_url)

    # Replace post screenshot
    if post_ss:
        replace_shape_with_image(slide, "ph_post_image", post_ss,
                                 horizontal_align="left")
        remove_text_shape_by_text(slide, "{POST_IMAGE}")
    else:
        # Keep the placeholder rect, just clean up
        remove_text_shape_by_text(slide, "{POST_IMAGE}")

    # Replace insight screenshot
    if insight_ss:
        replace_shape_with_image(slide, "ph_insight_image", insight_ss,
                                 horizontal_align="left")
        remove_text_shape_by_text(slide, "{INSIGHT_IMAGE}")
    else:
        remove_text_shape_by_text(slide, "{INSIGHT_IMAGE}")


def _fill_summary(slide, campaign_data: dict, campaign_name: str, campaign_month: str):
    """Fill the campaign summary slide."""
    xhs = campaign_data.get("xhs_kocs", [])
    tiktok = campaign_data.get("tiktok_kols", [])

    # XHS totals
    xhs_likes = total(xhs, "likes")
    xhs_comments = total(xhs, "comments")
    xhs_views = total(xhs, "views")
    xhs_saves = total(xhs, "saves")
    xhs_shares = total(xhs, "shares")
    xhs_posts = len(xhs)

    # TikTok totals
    tt_likes = total(tiktok, "likes")
    tt_comments = total(tiktok, "comments")
    tt_views = total(tiktok, "views")
    tt_saves = total(tiktok, "saves")
    tt_shares = total(tiktok, "shares")
    tt_posts = len(tiktok)

    # IG (combined from both platforms' cross-posts)
    ig_xhs = [k for k in xhs if k.get("ig_post_url")]
    ig_tt = [k for k in tiktok if k.get("ig_post_url")]
    ig_likes = sum(k.get("ig_likes", 0) for k in ig_xhs + ig_tt)
    ig_comments = sum(k.get("ig_comments", 0) for k in ig_xhs + ig_tt)
    ig_views = sum(k.get("ig_views", 0) for k in ig_xhs + ig_tt)
    ig_saves = sum(k.get("ig_saves", 0) for k in ig_xhs + ig_tt)
    ig_shares = sum(k.get("ig_shares", 0) for k in ig_xhs + ig_tt)
    ig_posts = len(ig_xhs) + len(ig_tt)

    # Totals
    total_likes = xhs_likes + tt_likes + ig_likes
    total_comments = xhs_comments + tt_comments + ig_comments
    total_views = xhs_views + tt_views + ig_views
    total_saves = xhs_saves + tt_saves + ig_saves
    total_shares = xhs_shares + tt_shares + ig_shares
    total_posts = xhs_posts + tt_posts + ig_posts
    total_engagement = total_likes + total_comments + total_saves + total_shares

    summary_text = (
        f"The campaign delivered positive results, garnering a total of {fmt(total_views)} views and "
        f"{fmt(total_engagement)} engagements (likes, comments, saves, and shares). "
        f"{campaign_name} {campaign_month} has successfully seeded a total of {total_posts} posts "
        f"across XHS, TikTok and Instagram."
    )

    replace_text_in_slide(slide, {
        # Tiktok row
        "{TIKTOK_POSTS}": fmt(tt_posts),
        "{TIKTOK_LIKES}": fmt(tt_likes),
        "{TIKTOK_COMMENTS}": fmt(tt_comments),
        "{TIKTOK_VIEWS}": fmt(tt_views),
        "{TIKTOK_SAVES}": fmt(tt_saves),
        "{TIKTOK_SHARES}": fmt(tt_shares),
        # XHS row
        "{XHS_POSTS}": fmt(xhs_posts),
        "{XHS_LIKES}": fmt(xhs_likes),
        "{XHS_COMMENTS}": fmt(xhs_comments),
        "{XHS_VIEWS}": fmt(xhs_views),
        "{XHS_SAVES}": fmt(xhs_saves),
        "{XHS_SHARES}": fmt(xhs_shares),
        # IG row
        "{IG_POSTS}": fmt(ig_posts),
        "{IG_LIKES}": fmt(ig_likes),
        "{IG_COMMENTS}": fmt(ig_comments),
        "{IG_VIEWS}": fmt(ig_views),
        "{IG_SAVES}": fmt(ig_saves),
        "{IG_SHARES}": fmt(ig_shares),
        # Total row
        "{TOTAL_POSTS}": fmt(total_posts),
        "{TOTAL_LIKES}": fmt(total_likes),
        "{TOTAL_COMMENTS}": fmt(total_comments),
        "{TOTAL_VIEWS}": fmt(total_views),
        "{TOTAL_SAVES}": fmt(total_saves),
        "{TOTAL_SHARES}": fmt(total_shares),
        # Summary paragraph
        "{SUMMARY_PARAGRAPH}": summary_text,
    })
